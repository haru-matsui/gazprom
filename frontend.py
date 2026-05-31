import streamlit as st
import requests
import folium
from streamlit_folium import st_folium
import base64
from io import BytesIO
import time
from PIL import Image
import json
import streamlit.components.v1 as components
import os
import re
import random
from pathlib import Path
from dotenv import load_dotenv
from playwright.sync_api import sync_playwright
from pptx import Presentation
from pptx.util import Inches
from pypdf import PdfReader, PdfWriter

# load .env from project root (if present)
load_dotenv()


def _safe_filename(value, fallback="presentation"):
        cleaned = re.sub(r'[<>:"/\\|?*\x00-\x1f]+', "_", str(value)).strip()
        cleaned = re.sub(r"\s+", " ", cleaned)
        return cleaned or fallback


def build_presentation_document(slide_css, slides_html, reset_script=""):
        return f"""<!doctype html>
<html lang="ru">
<head>
    <meta charset="utf-8">
    <meta name="viewport" content="width=device-width, initial-scale=1.0">
    {slide_css}
    <style>
        @media print {{
            html, body {{
                margin: 0;
                padding: 0;
                background: #ffffff;
                -webkit-print-color-adjust: exact;
                print-color-adjust: exact;
            }}
            .presentation-wrapper {{
                display: block !important;
                height: auto !important;
                overflow: visible !important;
                gap: 0 !important;
                padding: 0 !important;
                margin: 0 !important; /* Сбрасываем экранный margin-top: 20px, ломавший верстку */
                background: transparent !important;
            }}
            .slide {{
                display: block !important;
                break-inside: avoid !important;
                page-break-inside: avoid !important;
                page-break-after: avoid !important; /* Убираем всегда, так как слайды рендерятся по одному */
                break-after: avoid !important;      /* Предотвращает появление пустой/разорванной второй страницы */
                margin: 0 auto 0 auto !important;
                box-shadow: none !important;
                height: 562px !important;     /* Выравниваем высоту точно под размер страницы @page */
                max-height: 562px !important; /* Запрещаем контейнеру растягиваться */
                overflow: hidden !important;
                box-sizing: border-box !important;
            }}
            .slide-header {{
                padding: 18px 40px !important;
                font-size: 24px !important;
            }}
            .slide-body {{
                padding: 24px 34px !important;
                gap: 28px !important;
            }}
            .title-slide {{
                background: #d9d9d9 !important;
                color: #222 !important;
                justify-content: flex-start !important;
                align-items: flex-start !important;
                text-align: left !important;
                padding: 28px !important;
                height: 562px !important; /* Задаем жесткую высоту и для титульника */
                box-sizing: border-box !important;
            }}
            .title-slide h1 {{
                font-size: 34px !important;
                margin: 0 0 10px 0 !important;
                color: #222 !important;
            }}
            .title-slide h2 {{
                font-size: 24px !important;
                margin: 0 0 10px 0 !important; /* Сбрасываем дефолтные маргины браузера, раздувавшие слайд */
                color: #444 !important;
            }}
            .title-subtitle {{
                margin-top: 18px !important;
                padding-top: 0 !important;
                border-top: none !important;
                width: auto !important;
                max-width: 100% !important;
            }}
            .title-subtitle p {{
                font-size: 18px !important;
                color: #666 !important;
                margin: 0 !important;
            }}
            .slide-end {{
                display: none !important;
            }}
            .slide:last-child {{
                page-break-after: avoid !important;
                break-after: avoid !important;
            }}
            .slide-end {{
                color: #ffffff;
            }}
        }}
        @page {{
            size: 1000px 562px;
            margin: 0;
        }}
    </style>
</head>
<body>
{slides_html}
{reset_script}
</body>
</html>"""


def create_pdf_bytes(presentation_html, slide_css):
    with sync_playwright() as p:
        browser = p.chromium.launch(headless=True)
        page = browser.new_page(viewport={"width": 1400, "height": 900}, device_scale_factor=1)
        page.set_content(presentation_html, wait_until="networkidle")
        page.emulate_media(media="print")

        slides = page.locator(".slide:not(.slide-end)")
        slide_count = slides.count()
        writer = PdfWriter()

        for index in range(slide_count):
            slide_html = slides.nth(index).evaluate("element => element.outerHTML")
            single_slide_html = build_presentation_document(slide_css, f'<div class="presentation-wrapper">{slide_html}</div>')
            slide_page = browser.new_page(viewport={"width": 1400, "height": 900}, device_scale_factor=1)
            slide_page.set_content(single_slide_html, wait_until="networkidle")
            slide_page.emulate_media(media="print")
            slide_pdf = slide_page.pdf(
                print_background=True,
                prefer_css_page_size=True,
                scale=0.92,
            )
            slide_page.close()
            reader = PdfReader(BytesIO(slide_pdf))
            for pdf_page in reader.pages:
                writer.add_page(pdf_page)

        output = BytesIO()
        writer.write(output)
        browser.close()
        return output.getvalue()


def create_pptx_bytes(presentation_html):
        with sync_playwright() as p:
                browser = p.chromium.launch(headless=True)
                page = browser.new_page(viewport={"width": 1200, "height": 900}, device_scale_factor=1)
                page.set_content(presentation_html, wait_until="networkidle")
                page.emulate_media(media="screen")

                slides = page.locator(".slide:not(.slide-end)")
                slide_count = slides.count()
                presentation = Presentation()
                presentation.slide_width = Inches(13.333)
                presentation.slide_height = Inches(7.5)
                blank_layout = presentation.slide_layouts[6]

                for index in range(slide_count):
                        slide = presentation.slides.add_slide(blank_layout)
                        image_bytes = slides.nth(index).screenshot(type="png")
                        slide.shapes.add_picture(BytesIO(image_bytes), 0, 0, width=presentation.slide_width, height=presentation.slide_height)

                output = BytesIO()
                presentation.save(output)
                browser.close()
                return output.getvalue()
def generate_concept_board_openrouter(region_name, cultural_code):
    # API key must be provided via environment variable `OPENROUTER_API_KEY` or a .env file
    api_key = os.getenv("OPENROUTER_API_KEY")
    if not api_key:
        st.error("API key not found. Set OPENROUTER_API_KEY in your environment or in a .env file.")
        return None
    model_name = "bytedance-seed/seedream-4.5"
    seed_value = 777
    style = cultural_code.get('architecture_style', '')
    materials = cultural_code.get('materials', '')
    colors = ", ".join(cultural_code.get('color_profile', []))
    prompt = (
        f"An architectural **mood board layout and presentation panel** for a modern **sandwich-panel** factory concept in {region_name}. "
        f"The board must be a **collection of distinct, separate elements** arranged on a clean background. "
        f"Architectural style: {style}, design should reflect the local cultural code. "
        f"Materials to use: {materials}. "
        f"Color palette: {colors}. "
        f"**Do not generate a single, complete building render.** "
        f"Instead, the board must feature **multiple vignettes**: "
        "1. A **section view** or diagram of the facade showing traditional pattern integration. "
        "2. **Material swatches** (physical samples like wood, composite panels with hex codes). "
        "3. An **isometric view** of a building corner with texture overlays. "
        "4. A small **sketch or diagram** of the 'cultural code' concept. "
        "5. The color palette squares with hex codes. "
        f"Professional design presentation board, high resolution, clean russian typography, multiple viewports, concept-driven layout."
    )
    # 3. Собираем запрос
    headers = {
        "Authorization": f"Bearer {api_key}",
        "HTTP-Referer": "https://localhost",
        "X-Title": "Scoring App"
    }
    payload = {
        "model": model_name,
        "messages": [
            {"role": "user", "content": prompt}
        ],
        "modalities": ["image"],
        "seed": seed_value
    }
    response = requests.post(
        url="https://openrouter.ai/api/v1/chat/completions",
        headers=headers,
        json=payload
    )
    if response.status_code == 200:
        data = response.json()
        message = data["choices"][0]["message"]
        image_url_or_b64 = None
        if "images" in message and len(message["images"]) > 0:
            img_data = message["images"][0]
            if isinstance(img_data, dict) and "image_url" in img_data:
                image_url_or_b64 = img_data["image_url"]["url"]
            elif isinstance(img_data, str):
                image_url_or_b64 = img_data
        if not image_url_or_b64 and "content" in message:
            content = message["content"]
            if "data:image/" in content:
                start_idx = content.find("data:image/")
                end_chars = [content.find(")", start_idx), content.find('"', start_idx), content.find("'", start_idx)]
                valid_ends = [x for x in end_chars if x > start_idx]
                end_idx = min(valid_ends) if valid_ends else len(content)
                image_url_or_b64 = content[start_idx:end_idx]
        if image_url_or_b64 and "," in image_url_or_b64:
            base64_data = image_url_or_b64.split(",")[1]
            image_bytes = base64.b64decode(base64_data)
            return Image.open(BytesIO(image_bytes))
    else:
        st.error(f"Ошибка API: {response.status_code} - {response.text}")
    return None


def load_factory_collage_image(max_attempts=12, delay_seconds=0.25):
    collage_url = "http://localhost:8000/renders/factory_collage.png"
    for _ in range(max_attempts):
        try:
            response = requests.get(collage_url, timeout=3)
            if response.status_code == 200 and response.content:
                return Image.open(BytesIO(response.content))
        except Exception:
            pass
        time.sleep(delay_seconds)
    return None


def generate_collage_from_scene(three_js_data, viewpoints=None, out_dir="renders"):
    """
    Делает 4 снимка сцены из `three_scene.html` по заданным точкам на земле.
    Если точки не переданы, используются координаты, которые прислал пользователь.
    """
    viewpoints = viewpoints or [
        {"x": -60, "z": -10},
        {"x": 560, "z": -10},
        {"x": 250, "z": -280},
        {"x": 250, "z": 260},
    ]

    saved_paths = []
    out_path = Path(out_dir)

    try:
        with sync_playwright() as p:
            browser = p.chromium.launch(headless=True)
            page = browser.new_page(viewport={"width": 1200, "height": 800}, device_scale_factor=1)
            try:
                with open("three_scene.html", "r", encoding="utf-8") as f:
                    html_template = f.read()
            except FileNotFoundError:
                browser.close()
                return []

            html = html_template.replace("__DATA_PLACEHOLDER__", json.dumps(three_js_data))
            page.set_content(html, wait_until="networkidle")
            page.emulate_media(media="screen")
            page.wait_for_function("window.__threeSceneReady === true", timeout=30000)
            out_path.mkdir(exist_ok=True)

            for old_file in out_path.glob("factory_view_*.png"):
                try:
                    old_file.unlink()
                except OSError:
                    pass

            raw_points = [{"x": float(v.get("x", 0.0)), "z": float(v.get("z", 0.0))} for v in viewpoints[:4]]
            data_urls = page.evaluate("points => window.captureFactoryShotsFromPoints(points)", raw_points)
            if isinstance(data_urls, str):
                data_urls = [data_urls]

            for idx, data_url in enumerate(data_urls[:4], start=1):
                if "," in data_url:
                    _, encoded = data_url.split(",", 1)
                else:
                    encoded = data_url
                file_path = out_path / f"factory_view_{idx}.png"
                file_path.write_bytes(base64.b64decode(encoded))
                saved_paths.append(str(file_path))

            if saved_paths:
                try:
                    collage_size = 1536
                    thumb_size = collage_size // 2
                    positions = [(0, 0), (thumb_size, 0), (0, thumb_size), (thumb_size, thumb_size)]
                    collage = Image.new("RGB", (collage_size, collage_size), "white")
                    for idx, image_path in enumerate(saved_paths[:4]):
                        with Image.open(image_path) as source:
                            frame = source.convert("RGB").resize((thumb_size, thumb_size), Image.Resampling.LANCZOS)
                            collage.paste(frame, positions[idx])
                    collage.save(out_path / "factory_collage.png", format="PNG")
                except Exception:
                    pass

            browser.close()
        return saved_paths
    except Exception:
        return []
st.set_page_config(page_title="Подбор локации", layout="wide")
st.title("Выбор оптимального региона для строительства")
st.markdown("Внесите требуемые параметры для скоринга:")
if "concept_boards" not in st.session_state:
    st.session_state["concept_boards"] = {}
if "four_renders" not in st.session_state:
    st.session_state["four_renders"] = {}
if "3d_screenshots" not in st.session_state:
    st.session_state["3d_screenshots"] = {} # Перевели в словарь по индексу вкладки
if "final_renders" not in st.session_state:
    st.session_state["final_renders"] = {} # Перевели в словарь по индексу вкладки
if "factory_collage" not in st.session_state:
    st.session_state["factory_collage"] = None
col1, col2 = st.columns(2)
with col1:
    st.subheader("1. Производство и Экономика")
    production_volume = st.number_input("Объем производства (тыс. м²/год):", min_value=1, value=1000, step=50)
    employees = st.number_input("Количество сотрудников:", min_value=10, value=100, step=10)
    budget_mln = st.number_input("Бюджет (млн руб):", min_value=10, value=1500, step=100)
    st.subheader("2. Логистика")
    railway_options = {
        "Не имеет значения": "any",
        "Обязательно нужны": "required",
        "Не требуются (запрещено)": "forbidden"
    }
    railway_choice = st.selectbox("Потребность в Ж/Д путях", list(railway_options.keys()))
    railway_mode = railway_options[railway_choice]
    max_distance_to_highway = st.slider("Макс. удаленность от трассы (км)", min_value=0, max_value=100, value=10)
with col2:
    st.subheader("3. Архитектура и Благоустройство")
    architecture_priority = st.selectbox("Архитектурный приоритет",
                                         ["Экодизайн", "Техно-стиль", "Аутентичность региону"])
    landscaping_options = ["Сквер с фонтаном", "Беседки", "Сцена", "Пруд", "Арт-объект"]
    landscaping = st.multiselect("Элементы благоустройства", landscaping_options, default=["Пруд"])
    st.subheader("4. Социальные приоритеты")
    housing_percent = st.slider("Процент сотрудников с жильем (%)", min_value=0, max_value=100, value=10)
    housing_type = st.selectbox("Тип жилья", ["общежитие", "квартира"])
    kindergarten_places = st.number_input("Места в детском саду", min_value=0, value=50)
    sports_options = ["Уличные тренажёры", "Стадион", "Бассейн", "Спортзал", "Хоккейная коробка"]
    sports = st.multiselect("Спортивные объекты", sports_options, default=["Спортзал", "Стадион"])
if st.button("Рассчитать оптимальные регионы", type="primary"):
    payload = {
        "production_volume": production_volume,
        "employees": employees,
        "budget_mln": budget_mln,
        "railway_mode": railway_mode,
        "max_distance_to_highway": max_distance_to_highway,
        "architecture_priority": architecture_priority,
        "landscaping": landscaping,
        "housing": {
            "percent": housing_percent,
            "type": housing_type
        },
        "kindergarten_places": kindergarten_places,
        "sports": sports
    }
    with st.spinner("Анализ данных..."):
        try:
            response = requests.post("http://localhost:8000/api/score", json=payload)
            st.session_state["calculated"] = True
            st.session_state["top_regions"] = response.json().get("top_regions", [])
        except Exception as e:
            st.error(f"Ошибка при подключении к backend: {e}")
if st.session_state.get("calculated", False):
    top_regions = st.session_state["top_regions"]
    if not top_regions:
        st.warning("Нет регионов, удовлетворяющих жестким фильтрам.")
    else:
        if top_regions and "coordinates" in top_regions[0]["region"]:
            first_lat = top_regions[0]["region"]["coordinates"].get("lat", 54.0)
            first_lon = top_regions[0]["region"]["coordinates"].get("lon", 45.0)
            m = folium.Map(location=[first_lat, first_lon], zoom_start=6)
        else:
            m = folium.Map(location=[54.0, 45.0], zoom_start=5)
        from branca.element import Element
        m.get_root().html.add_child(
            Element("<style>.leaflet-control-attribution { display: none !important; }</style>"))
        tabs = st.tabs([f"{i + 1}. {r['region']['name']}" for i, r in enumerate(top_regions)])
        colors = ["red", "orange", "blue", "green", "purple"]
        for i, r in enumerate(top_regions):
            region_data = r["region"]
            est = r["estimate"]
            lat = region_data.get("coordinates", {}).get("lat", 54.0)
            lon = region_data.get("coordinates", {}).get("lon", 45.0)
            marker_color = colors[i] if i < len(colors) else "gray"
            icon_type = "star" if i == 0 else "info-sign"
            popup_html = f"""
            <div style="font-family: Arial, sans-serif; font-size:13px; line-height:1.2;">
              <b>{region_data.get('name', '')}</b> — Место: {i + 1} / Рейтинг: {r.get('score', '')}<br>
              <ul style="padding-left:18px; margin:6px 0 0 0; font-size:12px;">
                <li><strong>Льготы:</strong> {('ОЭЗ' if region_data.get('economics', {}).get('has_oez_benefits') else 'Нет')}</li>
                <li><strong>Сети:</strong> {('Газ' if region_data.get('infrastructure', {}).get('has_gas_in_promzone') else 'Нет')} / {region_data.get('infrastructure', {}).get('available_power_kva', '—')} кВА</li>
                <li><strong>Логистика:</strong> трасса {region_data.get('logistics', {}).get('distance_to_highway_km', '—')} км, Ж/Д: {('Есть' if region_data.get('logistics', {}).get('has_railway') else 'Нет')}</li>
                <li><strong>Экономика:</strong> тариф {region_data.get('economics', {}).get('energy_tariff_rub_kwh', '—')} руб/кВт·ч, з/п {region_data.get('economics', {}).get('average_salary_rub', '—')} руб</li>
                <li><strong>Социалка:</strong> индекс {region_data.get('social', {}).get('city_environment_index', '—')}, колледжи: {('Есть' if region_data.get('social', {}).get('has_profile_colleges') else 'Нет')}, сады {region_data.get('social', {}).get('kindergarten_occupancy_per_100_kids', '—')}/100</li>
              </ul>
            </div>
            """
            folium.Marker(
                location=[lat, lon],
                popup=folium.Popup(popup_html, max_width=320),
                tooltip=f"#{i + 1}: {region_data.get('name', '')} (Рейтинг: {r.get('score', '')})",
                icon=folium.Icon(color=marker_color, icon=icon_type)
            ).add_to(m)
            with tabs[i]:
                st.subheader(f"Итоговый рейтинг: {r['score']}")
                col_a, col_b = st.columns(2)
                with col_a:
                    st.markdown("#### Социальный паспорт региона")
                    st.write(f"- **Индекс городской среды:** {region_data['social']['city_environment_index']}")
                    st.write(
                        f"- **Обеспеченность детскими садами:** {region_data['social']['kindergarten_occupancy_per_100_kids']} мест на 100 детей")
                    st.write(
                        f"- **Профильные колледжи:** {'Есть' if region_data['social']['has_profile_colleges'] else 'Нет'}")
                    st.write(
                        f"- **Средняя аренда 1-комн. квартиры:** {region_data['social']['rent_1room_apartment_rub']} руб/мес")
                    st.markdown("#### Экономический блок")
                    st.write(
                        f"- **Налоговые льготы:** {'Есть' if region_data['economics']['has_oez_benefits'] else 'Нет'}")
                    # Показываем краткое описание преимуществ ОЭЗ/региона рядом с льготами
                    oez_adv = region_data.get('oez_advantage') or region_data.get('economics', {}).get('oez_advantage')
                    if oez_adv:
                        st.write(f"- **Преимущества региона / ОЭЗ:** {oez_adv}")
                    st.write(f"- **Энерготариф:** {region_data['economics']['energy_tariff_rub_kwh']} руб/кВт·ч")
                    st.write(f"- **Средняя зарплата:** {region_data['economics']['average_salary_rub']} руб/мес")
                    st.markdown("#### Рекомендации по удержанию персонала")
                    if region_data['social']['rent_1room_apartment_rub'] >= 25000:
                        st.write("- Рекомендуется строительство корпоративного жилья.")
                    else:
                        st.write("- Допускается использование рынка аренды жилья.")
                    if region_data['logistics']['distance_to_highway_km'] > 15:
                        st.write("- Требуется корпоративный автобус для сотрудников.")
                    else:
                        st.write("- Транспортная доступность соответствует нормативам.")
                    if region_data['social']['has_profile_colleges']:
                        st.write("- Рекомендуется сотрудничество с профильными колледжами.")
                with col_b:
                    st.markdown("#### Сетевой блок")
                    st.write(
                        f"- **Наличие газа:** {'Есть' if region_data['infrastructure']['has_gas_in_promzone'] else 'Нет'}")
                    st.write(f"- **Свободная мощность:** {region_data['infrastructure']['available_power_kva']} кВА")
                    st.write(
                        f"- **Стоимость подключения:** {region_data['infrastructure']['connection_cost_rub_kwt']} руб/кВт")
                    st.markdown("#### Логистика сырья")
                    st.write(
                        f"- **Расстояние до поставщика стали:** {region_data['logistics']['distance_to_steel_km']} км")
                    steel_name = region_data['logistics'].get('steel_supplier_name')
                    if steel_name:
                        st.write(f"  - Поставщик: {steel_name}")
                    st.write(
                        f"- **Расстояние до поставщика утеплителя:** {region_data['logistics']['distance_to_polyurethane_km']} км")
                    poly_name = region_data['logistics'].get('polyurethane_supplier_name')
                    if poly_name:
                        st.write(f"  - Поставщик: {poly_name}")
                    st.markdown("#### Предварительная смета строительства")
                    st.write(f"- **Площадь производственного цеха:** {est['shop_area']:,} м²")
                    st.write(f"- **Площадь склада:** {est['warehouse_area']:,} м²")
                    st.write(f"- **Площадь АБК:** {est['abk_area']:,} м²")
                    st.write(f"- **Площадь жилья:** {est['housing_area']:,} м²")
                    st.write(f"- **Площадь детского сада:** {est['kindergarten_area']:,} м²")
                    st.write("---")
                    st.write(f"- **CAPEX проекта:** {est['capex']:,} руб")
                    st.write(f"- **Годовой OPEX:** {est['opex']:,} руб/год")
                    st.markdown(f"**ИТОГОВАЯ СМЕТА ПРОЕКТА:** :green[**{est['total_cost']:,} руб**]")
                st.markdown("---")
                st.markdown("### Интерактивная 3D-визуализация площадки (Three.js)")
                st.write("Интерактивная трехмерная схема генерального плана на основе ваших параметров.")
                current_cultural_code = region_data.get('cultural_code', {})
                colors_from_code = current_cultural_code.get('color_profile', ["#555555", "#999999", "#cccccc"])
                three_js_data = {
                    "production_volume": production_volume,
                    "shop_area": est['shop_area'],
                    "warehouse_area": est['warehouse_area'],
                    "abk_area": est['abk_area'],
                    "housing_area": est['housing_area'],
                    "kindergarten_area": est['kindergarten_area'],
                    "landscaping": st.session_state.get("payload", {}).get('landscaping', []),
                    "sports": st.session_state.get("payload", {}).get('sports', []),
                    "color_palette": colors_from_code if colors_from_code else ["#555555", "#999999", "#cccccc"]
                }
                try:
                    with open("three_scene.html", "r", encoding="utf-8") as f:
                        html_template = f.read()
                    runtime_html_base = html_template.replace("__DATA_PLACEHOLDER__", json.dumps(three_js_data))
                    # site-only: button to (re)generate collage by reloading the embedded scene
                except FileNotFoundError:
                    st.error("Не удалось найти файл шаблона сцены `three_scene.html` в папке проекта.")
                generate_pressed = st.button("Сгенерировать коллаж", key=f"generate_collage_{i}")
                if generate_pressed:
                    # clear cached preview
                    st.session_state["factory_collage"] = None
                    with st.spinner("Сохраняем скриншоты локально..."):
                        try:
                            saved = generate_collage_from_scene(three_js_data)
                            st.session_state["factory_collage_paths"] = saved
                            if saved:
                                st.success(f"Сохранено {len(saved)} изображений в папке renders/")
                                try:
                                    st.image(saved[0], caption=os.path.basename(saved[0]), width=760)
                                except Exception:
                                    pass
                            else:
                                st.error("Не удалось сохранить скриншоты. Проверьте three_scene.html и логи Playwright.")
                        except Exception as e:
                            st.error(f"Ошибка сохранения скринов: {e}")
                components.html(runtime_html_base, height=550, scrolling=False)
                st.markdown("#### Коллаж завода")
                if st.session_state["factory_collage"] is None:
                    st.session_state["factory_collage"] = load_factory_collage_image()
                if st.session_state["factory_collage"] is not None:
                    st.image(st.session_state["factory_collage"], caption="Собранный коллаж завода", width=760)
                else:
                    st.info("Коллаж ещё сохраняется. Подождите секунду и обновите страницу.")
                st.markdown("---")
                st.markdown("#### Архитектурный концепт-борд (AI Генерация)")
                st.write(
                    "Создание визуализации на основе культурного кода региона (цвета, материалы, исторический стиль).")
                cultural_code = region_data.get("cultural_code")
                if cultural_code:
                    colors_html = "".join([
                        f'<div style="background-color:{c}; width:30px; height:30px; display:inline-block; border-radius:5px; margin-right:5px; border:1px solid #ccc;"></div>'
                        for c in cultural_code.get('color_profile', [])])
                    st.markdown(f"**Цветовой профиль:** {colors_html}", unsafe_allow_html=True)
                    if st.button(f"Сгенерировать концепт для {region_data['name']}", key=f"generate_btn_{i}"):
                        with st.spinner("Нейросеть анализирует стиль и генерирует рендер..."):
                            try:
                                img = generate_concept_board_openrouter(region_data['name'], cultural_code)
                                st.session_state["concept_boards"][i] = img
                            except Exception as e:
                                st.error(f"Ошибка: {e}")
                    if i in st.session_state["concept_boards"] and st.session_state["concept_boards"][i] is not None:
                        empty_left, img_container, empty_right = st.columns([1, 2, 1])
                        with img_container:
                            st.image(st.session_state["concept_boards"][i],
                                     caption=f"Концепт-борд: {region_data['name']}",
                                     width=500)
        st.markdown("### Карта рекомендованных площадок")
        st_folium(m, width=900, height=400)
        st.markdown("---")
        st.header("Презентация для администрации (Топ-1 регион)")
        top1 = top_regions[0]["region"]
        top1_est = top_regions[0]["estimate"]
        user_req = st.session_state.get("payload", {})
        slide_css = """
<style>
.presentation-wrapper {
    display: flex;
    flex-direction: column;
    align-items: center;
    gap: 80px;
    padding: 60px 0;
    background-color: #2b2b2b;
    border-radius: 15px;
    margin-top: 20px;
    height: 700px;
    overflow-y: scroll;
    scroll-snap-type: y mandatory;
    scroll-behavior: smooth;
}
.slide {
    width: 1000px;
    height: 562px;
    background: #e8e8e8;
    box-shadow: 0 12px 36px rgba(0,0,0,0.5);
    border-radius: 12px;
    display: flex;
    flex-direction: column;
    overflow: hidden;
    font-family: -apple-system, BlinkMacSystemFont, "Segoe UI", Roboto, Helvetica, Arial, sans-serif;
    position: relative;
    scroll-snap-align: center;
    flex: 0 0 auto;
}
.slide-header {
    background: #222222;
    color: white;
    padding: 24px 50px;
    font-size: 28px;
    font-weight: 600;
    display: flex;
    align-items: center;
}
.slide-body {
    padding: 40px 50px;
    flex: 1;
    display: flex;
    gap: 40px;
}
.slide-col {
    flex: 1;
    display: flex;
    flex-direction: column;
    gap: 16px;
}
.info-box {
    background: #d9d9d9;
    border-left: 5px solid #555555;
    padding: 18px 24px;
    border-radius: 6px;
}
.info-box.accent {
    border-left: 5px solid #333333;
    background: #cccccc;
}
.info-box h4 {
    margin: 0 0 8px 0;
    color: #444;
    font-size: 16px;
    font-weight: normal;
}
.info-box p {
    margin: 0;
    font-size: 26px;
    font-weight: bold;
    color: #222;
}
.info-box.accent p {
    color: #111;
}
.list-group {
    margin-top: 10px;
}
.list-item {
    font-size: 18px;
    color: #333;
    margin-bottom: 14px;
    line-height: 1.4;
}
.list-item strong {
    color: #111;
}
.title-slide {
    background: linear-gradient(135deg, #2b2b2b 0%, #1a1a1a 100%);
    color: white;
    justify-content: center;
    align-items: center;
    text-align: center;
    padding: 60px;
}
.title-slide h1 {
    font-size: 54px;
    margin-bottom: 24px;
    color: #ffffff;
    line-height: 1.2;
}
.title-slide h2 {
    font-size: 34px;
    color: #bbbbbb;
    font-weight: 300;
}
.section-title {
    font-size: 22px;
    color: #666666 !important;
    border-bottom: 2px solid #cccccc;
    padding-bottom: 10px;
    margin-bottom: 15px;
    margin-top: 0;
}
.presentation-wrapper::-webkit-scrollbar {
    width: 10px;
}
.presentation-wrapper::-webkit-scrollbar-track {
    background: #1a1a1a;
    border-radius: 10px;
}
.presentation-wrapper::-webkit-scrollbar-thumb {
    background: #555;
    border-radius: 10px;
}
.presentation-wrapper::-webkit-scrollbar-thumb:hover {
    background: #777;
}
.slide-end {
    background: linear-gradient(135deg, #141414 0%, #232323 100%);
    color: white;
}
.slide-end .slide-header {
    background: rgba(255, 255, 255, 0.08);
    color: #ffffff;
}
.slide-end .slide-body {
    align-items: center;
    justify-content: center;
    text-align: center;
}
.end-badge {
    display: inline-block;
    background: rgba(255, 255, 255, 0.12);
    border: 1px solid rgba(255, 255, 255, 0.22);
    border-radius: 999px;
    padding: 10px 18px;
    margin-bottom: 18px;
    font-size: 16px;
    letter-spacing: 0.08em;
    text-transform: uppercase;
}
.end-title {
    font-size: 42px;
    font-weight: 700;
    margin: 0 0 14px 0;
    color: #ffffff;
}
.end-text {
    font-size: 22px;
    color: #d8d8d8;
    line-height: 1.5;
    margin: 0;
}
</style>
"""
        st.markdown(slide_css, unsafe_allow_html=True)
        has_gas = "Да (промзона обеспечена)" if top1['infrastructure']['has_gas_in_promzone'] else "Нет"
        has_rail = "Есть доступ" if top1['logistics']['has_railway'] else "Отсутствует"
        has_oez = "ОЭЗ (Налоговые льготы)" if top1['economics']['has_oez_benefits'] else "Базовый режим"
        has_coll = "Присутствуют" if top1['social']['has_profile_colleges'] else "Отсутствуют"
        soc_area = top1_est['kindergarten_area'] + top1_est['housing_area']
        slides_html = f"""
<div class="presentation-wrapper">
<div class="slide title-slide">
<h1>Проект строительства<br>производственного комплекса</h1>
<h2>{top1['name']}</h2>
<div class="title-subtitle" style="margin-top: 50px; padding-top: 40px; border-top: 1px solid rgba(255,255,255,0.3); width: 60%;">
<p style="font-size: 24px; color: #a8c1e8; margin: 0;">Стратегия локализации на основе оценки региона</p>
</div>
</div>
<div class="slide">
<div class="slide-header">Параметры проекта и рабочие места</div>
<div class="slide-body">
<div class="slide-col">
<div class="info-box">
<h4>Общий budget проекта</h4>
<p>{user_req.get('budget_mln', 0)} млн руб.</p>
</div>
<div class="info-box accent">
<h4>Создаваемые рабочие места</h4>
<p>{user_req.get('employees', 0)} чел.</p>
</div>
<div class="info-box">
<h4>Плановый объем производства</h4>
<p>{user_req.get('production_volume', 0)} тыс. м²/год</p>
</div>
</div>
<div class="slide-col">
<h3 class="section-title">Масштаб строительства</h3>
<div class="list-group">
<div class="list-item"><strong>Площадь цеха:</strong> <span>{top1_est['shop_area']:,} м²</span></div>
<div class="list-item"><strong>Площадь складов:</strong> <span>{top1_est['warehouse_area']:,} м²</span></div>
<div class="list-item"><strong>CAPEX:</strong> <span style="color:#0056A4; font-weight:bold;">{top1_est['capex']:,} руб.</span></div>
</div>
<h3 class="section-title" style="margin-top: 20px;">Архитектура и среда</h3>
<div class="list-group">
<div class="list-item"><strong>Стиль:</strong> <span>{user_req.get('architecture_priority', 'Не указано')}</span></div>
<div class="list-item"><strong>Благоустройство:</strong> <span>{', '.join(user_req.get('landscaping', []))}</span></div>
</div>
</div>
</div>
</div>
<div class="slide">
<div class="slide-header">Соответствие нормативам и сети</div>
<div class="slide-body">
<div class="slide-col">
<h3 class="section-title">Сетевая инфраструктура</h3>
<div class="info-box accent">
<h4>Свободная электр. мощность</h4>
<p>{top1['infrastructure']['available_power_kva']} кВА</p>
</div>
<div class="list-group">
<div class="list-item"><strong>Газификация:</strong> <span>{has_gas}</span></div>
<div class="list-item"><strong>Стоимость тех. присоединения:</strong> <span>{top1['infrastructure']['connection_cost_rub_kwt']} руб/кВт</span></div>
<div class="list-item"><strong>Удалённость от подстанции:</strong> <span>{top1['infrastructure']['distance_to_substation_km']} км</span></div>
</div>
</div>
<div class="slide-col">
<h3 class="section-title">Логистические нормативы</h3>
<div class="list-group" style="background: #f8f9fa; padding: 20px; border-radius: 8px;">
<div class="list-item">🚂 <strong>Ж/Д ветка:</strong> <span>{has_rail}</span></div>
<div class="list-item">🛣 <strong>Расстояние до трассы:</strong> <span>{top1['logistics']['distance_to_highway_km']} км</span><br>
<em style="font-size: 14px; color: #666;">(норматив заявки: до {user_req.get('max_distance_to_highway', 0)} км)</em>
</div>
<div class="list-item">🏭 <strong>Сырьё (сталь):</strong> <span>{top1['logistics']['distance_to_steel_km']} км</span></div>
<div class="list-item">📦 <strong>Сырьё (утеплитель):</strong> <span>{top1['logistics']['distance_to_polyurethane_km']} км</span></div>
</div>
</div>
</div>
</div>
<div class="slide">
<div class="slide-header">Экономические и социальные выгоды</div>
<div class="slide-body">
<div class="slide-col">
<h3 class="section-title">Экономика проекта</h3>
<div class="info-box">
<h4>Статус территории (льготы)</h4>
<p style="font-size: 20px;">{has_oez}</p>
</div>
<div class="list-group">
<div class="list-item"><strong>Сниженный энерготариф:</strong> <span>{top1['economics']['energy_tariff_rub_kwh']} руб/кВт·ч</span></div>
<div class="list-item"><strong>Средняя зарплата по региону:</strong> <span>{top1['economics']['average_salary_rub']:,} руб/мес</span></div>
<div class="list-item"><strong>Годовой OPEX:</strong> <span>{top1_est['opex']:,} руб.</span></div>
</div>
</div>
<div class="slide-col">
<h3 class="section-title">Социальное развитие региона</h3>
<div class="list-group">
<div class="list-item"><strong>Индекс городской среды:</strong> <span style="background:#eaf1ed; padding:4px 10px; border-radius:4px; font-weight:bold; color: #111;">{top1['social']['city_environment_index']} баллов</span></div>
<div class="list-item"><strong>Мест в детсадах:</strong> <span>{top1['social']['kindergarten_occupancy_per_100_kids']} (на 100 детей)</span></div>
<div class="list-item"><strong>Профильные колледжи:</strong> <span>{has_coll}</span></div>
</div>
<div class="info-box accent" style="margin-top: 15px;">
<h4>Собственные соц. объекты (сад, жилье)</h4>
<p style="font-size: 22px;">{"Не требуются" if soc_area == 0 else f"Площадь: {soc_area:,} м²"}</p>
</div>
</div>
</div>
</div>
<div class="slide slide-end">
<div class="slide-header">Конец презентации</div>
<div class="slide-body">
<div>
<div class="end-badge">Финальный слайд</div>
<div class="end-title">Спасибо за внимание</div>
<p class="end-text">Презентация завершена. Ниже можно скачать PDF и PPTX</p>
</div>
</div>
</div>
</div>
"""
        reset_script = '''
<script>
    setTimeout(function(){
        try{
            var el = document.querySelector('.presentation-wrapper');
            if(el){ el.scrollTo({ top: 0, left: 0, behavior: 'instant' }); }
        }catch(e){}
    }, 50);
</script>
'''
        st.markdown(slides_html + reset_script, unsafe_allow_html=True)
        presentation_html = build_presentation_document(slide_css, slides_html, reset_script)
        presentation_name = _safe_filename(top1['name'])
        export_key = f"{presentation_name}:{top1_est['total_cost']}"
        if st.session_state.get("presentation_export_key") != export_key:
            with st.spinner("Подготавливаем PDF и PPTX..."):
                try:
                    st.session_state["presentation_pdf_bytes"] = create_pdf_bytes(presentation_html, slide_css)
                    st.session_state["presentation_pptx_bytes"] = create_pptx_bytes(presentation_html)
                    st.session_state["presentation_export_key"] = export_key
                except Exception as export_error:
                    st.session_state["presentation_pdf_bytes"] = b""
                    st.session_state["presentation_pptx_bytes"] = b""
                    st.error(f"Не удалось подготовить файлы презентации: {export_error}")

        st.markdown("#### Скачать презентацию")
        download_left, download_right = st.columns(2)
        with download_left:
            st.download_button(
                label="Скачать PPTX",
                data=st.session_state.get("presentation_pptx_bytes", b""),
                file_name=f"{presentation_name}.pptx",
                mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                use_container_width=True,
                disabled=not st.session_state.get("presentation_pptx_bytes")
            )
        with download_right:
            st.download_button(
                label="Скачать PDF",
                data=st.session_state.get("presentation_pdf_bytes", b""),
                file_name=f"{presentation_name}.pdf",
                mime="application/pdf",
                use_container_width=True,
                disabled=not st.session_state.get("presentation_pdf_bytes")
            )